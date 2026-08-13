from pathlib import Path
from typing import Union, Dict, Any
from app.loaders.base import BaseLoader
from app.models.internal import Document
from app.core.errors import InvalidFileError
from app.core.logging import logger

try:
    import pptx
except ImportError:
    pptx = None


class PPTXLoader(BaseLoader):
    """PPTX presentation loader extracting slide titles, shape text, and slide notes per slide."""

    def load(self, source: Union[str, Path], metadata: Dict[str, Any] = None) -> Document:
        file_path = Path(source)
        if not file_path.exists():
            raise InvalidFileError(f"PPTX file not found: {source}")

        if pptx is None:
            raise InvalidFileError("python-pptx package is not installed.")

        try:
            prs = pptx.Presentation(str(file_path))
            slides_text = []

            for idx, slide in enumerate(prs.slides):
                slide_num = idx + 1
                slide_lines = [f"[Slide {slide_num}]"]

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text.strip()
                        if text:
                            slide_lines.append(text)
                    elif shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_vals = [cell.text.strip() for cell in row.cells]
                            if any(row_vals):
                                slide_lines.append(" | ".join(row_vals))

                # Slide notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_lines.append(f"Speaker Notes: {notes}")

                slides_text.append("\n".join(slide_lines))

            full_content = "\n\n".join(slides_text)

            meta = {
                "source_type": "pptx",
                "file_name": file_path.name,
                "slide_count": len(prs.slides),
                "file_size": file_path.stat().st_size,
                **(metadata or {})
            }

            return Document(
                source_type="pptx",
                source_name=file_path.name,
                source_uri=str(file_path.absolute()),
                content=full_content,
                metadata=meta
            )
        except Exception as e:
            logger.error(f"Failed to load PPTX file {source}: {e}")
            raise InvalidFileError(f"Failed to read PPTX file {source}: {str(e)}")
