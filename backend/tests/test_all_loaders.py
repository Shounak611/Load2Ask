import io
import json
import pytest
from pathlib import Path
import pypdf
from PIL import Image

from app.loaders.factory import LoaderFactory
from app.loaders.txt_loader import TextLoader
from app.loaders.pdf_loader import PDFLoader
from app.loaders.image_loader import ImageLoader
from app.loaders.web_loader import WebLoader, validate_url_ssrf, SSRFProtectionError
from app.loaders.docx_loader import DOCXLoader
from app.loaders.pptx_loader import PPTXLoader
from app.loaders.csv_loader import CSVLoader
from app.loaders.xlsx_loader import XLSXLoader
from app.loaders.json_loader import JSONLoader
from app.loaders.markdown_loader import MarkdownLoader
from app.loaders.html_loader import HTMLLoader


def test_txt_loader(tmp_path):
    txt_file = tmp_path / "sample.txt"
    txt_file.write_text("Hello World! TXT Loader test content.", encoding="utf-8")
    loader = TextLoader()
    doc = loader.load(txt_file)
    assert doc.source_type == "txt"
    assert "Hello World!" in doc.content


def test_pdf_loader(tmp_path):
    pdf_file = tmp_path / "sample.pdf"
    writer = pypdf.PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with pdf_file.open("wb") as f:
        writer.write(f)

    loader = PDFLoader()
    doc = loader.load(pdf_file)
    assert doc.source_type == "pdf"
    assert doc.metadata["num_pages"] == 1


def test_image_loader(tmp_path):
    img_file = tmp_path / "test.png"
    img = Image.new("RGB", (100, 50), color="blue")
    img.save(img_file)

    loader = ImageLoader()
    doc = loader.load(img_file)
    assert doc.source_type == "image"
    assert doc.metadata["width"] == 100
    assert doc.metadata["height"] == 50


def test_docx_loader(tmp_path):
    import docx
    docx_file = tmp_path / "test.docx"
    doc_obj = docx.Document()
    doc_obj.add_heading("Docx Title", level=1)
    doc_obj.add_paragraph("Paragraph inside docx document.")
    doc_obj.save(str(docx_file))

    loader = DOCXLoader()
    doc = loader.load(docx_file)
    assert doc.source_type == "docx"
    assert "Docx Title" in doc.content
    assert "Paragraph inside docx document." in doc.content


def test_pptx_loader(tmp_path):
    import pptx
    pptx_file = tmp_path / "test.pptx"
    prs = pptx.Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    tx_box = slide.shapes.add_textbox(0, 0, 100, 100)
    tx_box.text_frame.text = "Slide 1 Content Text"
    prs.save(str(pptx_file))

    loader = PPTXLoader()
    doc = loader.load(pptx_file)
    assert doc.source_type == "pptx"
    assert "Slide 1 Content Text" in doc.content


def test_csv_loader(tmp_path):
    csv_file = tmp_path / "test.csv"
    csv_file.write_text("name,age,city\nAlice,30,New York\nBob,25,London", encoding="utf-8")
    loader = CSVLoader()
    doc = loader.load(csv_file)
    assert doc.source_type == "csv"
    assert "Alice" in doc.content
    assert "New York" in doc.content


def test_xlsx_loader(tmp_path):
    import openpyxl
    xlsx_file = tmp_path / "test.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "TestSheet"
    ws.append(["Header1", "Header2"])
    ws.append(["Value1", "Value2"])
    wb.save(str(xlsx_file))

    loader = XLSXLoader()
    doc = loader.load(xlsx_file)
    assert doc.source_type == "xlsx"
    assert "TestSheet" in doc.content
    assert "Value1" in doc.content


def test_json_loader(tmp_path):
    json_file = tmp_path / "test.json"
    data = {"project": "Load2Ask", "version": 2, "features": ["Multi-source", "RAG"]}
    json_file.write_text(json.dumps(data), encoding="utf-8")

    loader = JSONLoader()
    doc = loader.load(json_file)
    assert doc.source_type == "json"
    assert "Load2Ask" in doc.content
    assert "top_level_keys" in doc.metadata


def test_markdown_loader(tmp_path):
    md_file = tmp_path / "test.md"
    md_file.write_text("# Load2Ask RAG\n\n## Section 1\nContent of markdown file.", encoding="utf-8")

    loader = MarkdownLoader()
    doc = loader.load(md_file)
    assert doc.source_type == "markdown"
    assert "# Load2Ask RAG" in doc.content
    assert any("Section 1" in h for h in doc.metadata["headings"])



def test_html_loader(tmp_path):
    html_file = tmp_path / "test.html"
    html_file.write_text(
        "<html><head><title>HTML Test Page</title></head><body><h1>Main Title</h1><p>Body HTML paragraph.</p></body></html>",
        encoding="utf-8"
    )

    loader = HTMLLoader()
    doc = loader.load(html_file)
    assert doc.source_type == "html"
    assert "HTML Test Page" in doc.content
    assert "Body HTML paragraph." in doc.content


def test_ssrf_protection():
    # Test valid external URL passes validation
    assert validate_url_ssrf("https://example.com/page") == "https://example.com/page"

    # Test loopback and private IP rejection
    with pytest.raises(SSRFProtectionError):
        validate_url_ssrf("http://localhost:8000/admin")

    with pytest.raises(SSRFProtectionError):
        validate_url_ssrf("http://127.0.0.1:5432")

    with pytest.raises(SSRFProtectionError):
        validate_url_ssrf("http://192.168.1.1/secret")

    with pytest.raises(SSRFProtectionError):
        validate_url_ssrf("http://169.254.169.254/latest/meta-data")

    with pytest.raises(SSRFProtectionError):
        validate_url_ssrf("file:///etc/passwd")
